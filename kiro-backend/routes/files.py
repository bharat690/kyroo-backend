from fastapi import APIRouter, HTTPException
from pydantic import BaseModel
import os
import uuid
from pathlib import Path

router = APIRouter(prefix="/files", tags=["files"])

TEMP_DIR = "/tmp/kiro_files"
os.makedirs(TEMP_DIR, exist_ok=True)

class ConvertRequest(BaseModel):
    file_base64: str
    filename: str
    convert_to: str
    user_id: str

def get_extension(filename: str) -> str:
    return Path(filename).suffix.lower()

def pdf_to_word(input_path: str, output_path: str):
    from pdf2docx import Converter
    cv = Converter(input_path)
    cv.convert(output_path, start=0, end=None)
    cv.close()

def word_to_pdf(input_path: str, output_path: str):
    from docx import Document
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter
    
    doc = Document(input_path)
    c = canvas.Canvas(output_path, pagesize=letter)
    y = 750
    
    for para in doc.paragraphs:
        if para.text:
            c.drawString(50, y, para.text[:90])
            y -= 20
            if y < 50:
                c.showPage()
                y = 750
    c.save()

def ppt_to_pdf(input_path: str, output_path: str):
    from pptx import Presentation
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter
    
    prs = Presentation(input_path)
    c = canvas.Canvas(output_path, pagesize=letter)
    
    for i, slide in enumerate(prs.slides):
        c.drawString(50, 750, f"Slide {i + 1}")
        y = 700
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                c.drawString(50, y, shape.text[:90])
                y -= 25
                if y < 50:
                    break
        c.showPage()
    c.save()

def excel_to_pdf(input_path: str, output_path: str):
    import openpyxl
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter
    
    wb = openpyxl.load_workbook(input_path)
    c = canvas.Canvas(output_path, pagesize=letter)
    
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        c.drawString(50, 750, f"Sheet: {sheet}")
        y = 700
        for row in ws.iter_rows(values_only=True):
            row_text = " | ".join([str(cell) if cell else "" for cell in row])
            c.drawString(50, y, row_text[:90])
            y -= 20
            if y < 50:
                c.showPage()
                y = 750
        c.showPage()
    c.save()

def csv_to_excel(input_path: str, output_path: str):
    import openpyxl
    import csv
    
    wb = openpyxl.Workbook()
    ws = wb.active
    
    with open(input_path, 'r') as f:
        reader = csv.reader(f)
        for row in reader:
            ws.append(row)
    
    wb.save(output_path)

def images_to_pdf(image_paths: list, output_path: str):
    from PIL import Image
    
    images = []
    for path in image_paths:
        img = Image.open(path)
        if img.mode != 'RGB':
            img = img.convert('RGB')
        images.append(img)
    
    if images:
        images[0].save(
            output_path,
            save_all=True,
            append_images=images[1:]
        )

def read_pdf_text(input_path: str) -> str:
    import PyPDF2
    text = ""
    with open(input_path, 'rb') as f:
        reader = PyPDF2.PdfReader(f)
        for page in reader.pages:
            text += page.extract_text() + "\n"
    return text

@router.get("/supported")
async def supported_conversions():
    return {
        "conversions": [
            {"from": "PDF", "to": "Word (.docx)"},
            {"from": "PDF", "to": "Text (.txt)"},
            {"from": "Word (.docx)", "to": "PDF"},
            {"from": "PowerPoint (.pptx)", "to": "PDF"},
            {"from": "Excel (.xlsx)", "to": "PDF"},
            {"from": "CSV", "to": "Excel (.xlsx)"},
            {"from": "Image (jpg/png)", "to": "PDF"},
        ],
        "coming_soon": [
            {"from": "PDF", "to": "PowerPoint"},
            {"from": "PDF", "to": "Excel"},
            {"from": "Word", "to": "PowerPoint"},
        ],
        "status": "ready"
    }

@router.post("/convert")
async def convert_file(req: ConvertRequest):
    import base64
    
    try:
        file_data = base64.b64decode(req.file_base64)
    except:
        raise HTTPException(status_code=400, detail="Invalid file data")
    
    input_ext = get_extension(req.filename)
    convert_to = req.convert_to.lower().strip(".")
    
    input_id = str(uuid.uuid4())
    input_path = f"{TEMP_DIR}/{input_id}{input_ext}"
    output_path = f"{TEMP_DIR}/{input_id}.{convert_to}"
    
    with open(input_path, 'wb') as f:
        f.write(file_data)
    
    try:
        if input_ext == ".pdf" and convert_to == "docx":
            pdf_to_word(input_path, output_path)
        elif input_ext in [".docx", ".doc"] and convert_to == "pdf":
            word_to_pdf(input_path, output_path)
        elif input_ext in [".pptx", ".ppt"] and convert_to == "pdf":
            ppt_to_pdf(input_path, output_path)
        elif input_ext in [".xlsx", ".xls"] and convert_to == "pdf":
            excel_to_pdf(input_path, output_path)
        elif input_ext == ".csv" and convert_to in ["xlsx", "excel"]:
            output_path = f"{TEMP_DIR}/{input_id}.xlsx"
            csv_to_excel(input_path, output_path)
        elif input_ext in [".jpg", ".jpeg", ".png"] and convert_to == "pdf":
            images_to_pdf([input_path], output_path)
        elif input_ext == ".pdf" and convert_to == "txt":
            text = read_pdf_text(input_path)
            with open(output_path, 'w') as f:
                f.write(text)
        else:
            raise HTTPException(
                status_code=400,
                detail=f"ye conversion abhi supported nahi hai"
            )
        
        with open(output_path, 'rb') as f:
            output_data = base64.b64encode(f.read()).decode()
        
        output_filename = Path(req.filename).stem + f".{convert_to}"
        
        os.remove(input_path)
        os.remove(output_path)
        
        return {
            "file_base64": output_data,
            "filename": output_filename,
            "original": req.filename,
            "converted_to": convert_to,
            "status": "success",
            "message": f"ho gaya! {req.filename} ko {convert_to} mein convert kar diya 🎉"
        }
    
    except HTTPException:
        raise
    except Exception as e:
        if os.path.exists(input_path):
            os.remove(input_path)
        raise HTTPException(status_code=500, detail=f"conversion failed: {str(e)}")

@router.post("/read-pdf")
async def read_pdf_content(req: ConvertRequest):
    import base64
    
    try:
        file_data = base64.b64decode(req.file_base64)
    except:
        raise HTTPException(status_code=400, detail="Invalid file data")
    
    input_id = str(uuid.uuid4())
    input_path = f"{TEMP_DIR}/{input_id}.pdf"
    
    with open(input_path, 'wb') as f:
        f.write(file_data)
    
    try:
        text = read_pdf_text(input_path)
        os.remove(input_path)
        
        return {
            "text": text,
            "characters": len(text),
            "status": "success"
        }
    except Exception as e:
        if os.path.exists(input_path):
            os.remove(input_path)
        raise HTTPException(status_code=500, detail=str(e))